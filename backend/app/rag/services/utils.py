"""
Optimized Utility Functions for RAG System - SPEED OPTIMIZED (<60s)
- Parallel processing for document chunking
- Fast text processing and cleaning
- Optimized embedding generation
- Intelligent chunking with speed focus
"""

import re
import logging
import asyncio
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor
from typing import List, Dict, Any, Optional
import time
from functools import lru_cache
import hashlib

from app.rag.config import RAGConfig as Config

logger = logging.getLogger(__name__)

# Thread pools for parallel processing
chunking_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_CHUNKING)
embedding_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_EMBEDDINGS)
answer_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_ANSWERS)

def clean_text(text: str) -> str:
    """Fast text cleaning for speed optimization."""
    if not text:
        return ""
    
    # Fast cleaning operations
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces
    # Keep more characters including colons, dashes, and other important punctuation
    text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\[\]\{\}\+\=\*\/\@\#\$\%\&\*]', '', text)  # Keep more special chars
    text = text.strip()
    
    return text

def clean_text_fast(text: str) -> str:
    """Ultra-fast text cleaning for large documents."""
    if not text:
        return ""
    
    # Minimal cleaning for speed
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces
    text = text.strip()
    
    return text

def clean_text_for_images(text: str) -> str:
    """Specialized text cleaning for images that preserves mathematical expressions."""
    if not text:
        return ""
    
    # Preserve mathematical expressions by keeping +, =, -, *, /, (, ), [, ], {, }
    # Also preserve numbers and basic punctuation
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with single space
    text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\[\]\{\}\+\=\*\/]', '', text)  # Keep math operators
    text = text.strip()
    
    return text

def extract_text_from_file(file) -> str:
    """Fast text extraction with parallel processing."""
    try:
        filename = file.filename.lower()
        
        if filename.endswith('.zip'):
            return extract_text_from_zip_fast(file)
        elif filename.endswith('.pdf'):
            return extract_text_from_pdf_fast(file)
        elif filename.endswith(('.docx', '.doc')):
            return extract_text_from_docx_fast(file)
        elif filename.endswith('.pptx'):
            return extract_text_from_pptx_fast(file)
        elif filename.endswith(('.xlsx', '.xls')):
            return extract_text_from_excel_fast(file)
        elif filename.endswith('.csv'):
            return extract_text_from_csv_fast(file)
        elif filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
            return extract_text_from_image_fast(file)
        else:
            # For other file types, read as text
            content = file.read()
            return content.decode('utf-8', errors='ignore')
            
    except Exception as e:
        logger.error(f"Error extracting text from file: {e}")
        return ""

def extract_text_from_pdf_fast(file) -> str:
    """Ultra-fast PDF text extraction optimized for large files (<30s target)."""
    try:
        from pypdf import PdfReader
        
        # Read PDF in memory
        pdf_reader = PdfReader(file)
        total_pages = len(pdf_reader.pages)
        
        # For large PDFs (>100 pages), use intelligent sampling
        if total_pages > 100:
            return extract_text_from_large_pdf_fast(pdf_reader, total_pages)
        else:
            return extract_text_from_small_pdf_fast(pdf_reader, total_pages)
        
    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        return ""

def extract_text_from_large_pdf_fast(pdf_reader, total_pages: int) -> str:
    """Ultra-fast extraction for large PDFs using intelligent sampling."""
    # Calculate sampling strategy for 30-second target
    max_pages_to_process = min(50, total_pages // 4)  # Process max 50 pages
    
    # Sample pages intelligently: first 10, last 10, and every nth page in between
    pages_to_extract = []
    
    # Always include first 10 pages (usually introduction, table of contents)
    pages_to_extract.extend(range(min(10, total_pages)))
    
    # Always include last 10 pages (usually conclusion, references)
    if total_pages > 20:
        pages_to_extract.extend(range(max(10, total_pages - 10), total_pages))
    
    # Sample middle pages evenly
    if total_pages > 20:
        middle_start = 10
        middle_end = total_pages - 10
        step = max(1, (middle_end - middle_start) // (max_pages_to_process - 20))
        
        for i in range(middle_start, middle_end, step):
            if len(pages_to_extract) < max_pages_to_process:
                pages_to_extract.append(i)
    
    def extract_sampled_page_text(page_index):
        try:
            return pdf_reader.pages[page_index].extract_text()
        except:
            return ""
    
    # Use more workers for faster processing
    with ThreadPoolExecutor(max_workers=min(16, len(pages_to_extract))) as executor:
        texts = list(executor.map(extract_sampled_page_text, pages_to_extract))
    
    # Combine all texts
    full_text = " ".join(texts)
    return clean_text_fast(full_text)

def extract_text_from_small_pdf_fast(pdf_reader, total_pages: int) -> str:
    """Fast extraction for small PDFs (≤100 pages)."""
    def extract_page_text(page):
        try:
            return page.extract_text()
        except:
            return ""
    
    # Use thread pool for parallel extraction
    with ThreadPoolExecutor(max_workers=min(8, total_pages)) as executor:
        texts = list(executor.map(extract_page_text, pdf_reader.pages))
    
    # Combine all texts
    full_text = " ".join(texts)
    return clean_text(full_text)

def extract_text_from_docx_fast(file) -> str:
    """Fast DOCX text extraction."""
    try:
        import docx2txt
        text = docx2txt.process(file)
        cleaned_text = clean_text(text)
        return cleaned_text
    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        return ""

def chunk_text_advanced(text: str, chunk_size: Optional[int] = None, overlap: Optional[int] = None) -> List[str]:
    """Ultra-fast chunking optimized for large documents (<30s target)."""
    if not text:
        return []
    
    # Use config values if not provided
    chunk_size = chunk_size or Config.CHUNK_SIZE
    overlap = overlap or Config.CHUNK_OVERLAP
    
    # For very large texts (>500KB), use ultra-fast chunking
    if len(text) > 500000:  # 500KB threshold
        return chunk_text_ultra_fast(text, chunk_size, overlap)
    
    # Special handling for Excel data to preserve all rows for numerical comparisons
    if "SHEET:" in text and "Row" in text:
        return chunk_excel_data_specialized(text, chunk_size, overlap)
    
    # Special handling for PPTX data to prioritize slide content over notes
    if "SLIDE" in text and ("NOTES:" in text or "Power BI" in text):
        return chunk_pptx_data_specialized(text, chunk_size, overlap)
    
    # Fast chunking for speed
    if Config.ENABLE_FAST_CHUNKING:
        return chunk_text_parallel(text, chunk_size, overlap)
    else:
        return chunk_text_standard(text, chunk_size, overlap)

def chunk_excel_data_specialized(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Specialized chunking for Excel data to preserve all rows for numerical comparisons."""
    if not text:
        return []
    
    chunks = []
    lines = text.split('\n')
    
    # Find sheet headers and data rows
    sheet_headers = []
    data_rows = []
    
    for line in lines:
        if line.startswith('SHEET:') or line.startswith('  HEADERS:'):
            sheet_headers.append(line)
        elif line.startswith('  Row'):
            data_rows.append(line)
    
    # Create chunks that preserve complete data for numerical comparisons
    current_chunk = []
    current_length = 0
    
    # Add sheet headers to first chunk
    if sheet_headers:
        current_chunk.extend(sheet_headers)
        current_length += sum(len(header) for header in sheet_headers)
    
    # Process data rows
    for row in data_rows:
        row_length = len(row)
        
        # If adding this row would exceed chunk size, save current chunk and start new one
        if current_length + row_length > chunk_size and current_chunk:
            chunk_text = '\n'.join(current_chunk)
            if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
                chunks.append(chunk_text)
            
            # Start new chunk with overlap (keep last few rows)
            overlap_rows = current_chunk[-3:] if len(current_chunk) >= 3 else current_chunk
            current_chunk = overlap_rows + [row]
            current_length = sum(len(r) for r in current_chunk)
        else:
            current_chunk.append(row)
            current_length += row_length
    
    # Add final chunk
    if current_chunk:
        chunk_text = '\n'.join(current_chunk)
        if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
            chunks.append(chunk_text)
    
    # For Excel data, we want to ensure ALL rows are included for numerical comparisons
    # Don't limit chunks for Excel files to preserve complete data
    return chunks

def chunk_pptx_data_specialized(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Specialized chunking for PPTX data to prioritize slide content over notes."""
    if not text:
        return []
    
    chunks = []
    lines = text.split('\n')
    
    # Separate slide content from notes
    slide_content = []
    notes_content = []
    
    current_section = None
    for line in lines:
        if line.startswith('SLIDE') and 'NOTES:' in line:
            current_section = 'notes'
        elif line.startswith('SLIDE') and not 'NOTES:' in line:
            current_section = 'slide'
        
        if current_section == 'slide':
            slide_content.append(line)
        elif current_section == 'notes':
            notes_content.append(line)
        elif current_section is None:
            slide_content.append(line)
    
    # Temporarily include all notes to see what content exists
    filtered_notes = notes_content
    
    # Prioritize slide content, then filtered notes
    prioritized_content = slide_content + filtered_notes
    
    # Create chunks from prioritized content
    current_chunk = []
    current_length = 0
    
    for line in prioritized_content:
        line_length = len(line)
        
        # If adding this line would exceed chunk size, save current chunk and start new one
        if current_length + line_length > chunk_size and current_chunk:
            chunk_text = '\n'.join(current_chunk)
            if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
                chunks.append(chunk_text)
            
            # Start new chunk with overlap
            overlap_lines = current_chunk[-2:] if len(current_chunk) >= 2 else current_chunk
            current_chunk = overlap_lines + [line]
            current_length = sum(len(l) for l in current_chunk)
        else:
            current_chunk.append(line)
            current_length += line_length
    
    # Add final chunk
    if current_chunk:
        chunk_text = '\n'.join(current_chunk)
        if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
            chunks.append(chunk_text)
    
    # Don't limit chunks for PPTX to preserve all relevant content
    return chunks

def chunk_text_ultra_fast(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Ultra-fast chunking for very large documents."""
    if not text:
        return []
    
    # Simple word-based chunking for speed
    words = text.split()
    chunks = []
    
    # Calculate target chunks for large documents (limit to 100 chunks max)
    max_chunks = min(100, len(text) // chunk_size)
    words_per_chunk = max(1, len(words) // max_chunks)
    
    for i in range(0, len(words), words_per_chunk):
        chunk_words = words[i:i + words_per_chunk]
        chunk_text = " ".join(chunk_words)
        
        if len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
            chunks.append(chunk_text)
        
        # Limit chunks for speed
        if len(chunks) >= max_chunks:
            break
    
    return chunks

def chunk_text_parallel(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Parallel text chunking for speed optimization."""
    try:
        # Split text into sentences first
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        # Process chunks in parallel
        def create_chunk(sentence_batch):
            chunk_text = " ".join(sentence_batch)
            if len(chunk_text) > chunk_size:
                # If chunk is too large, split by words
                words = chunk_text.split()
                chunk_text = " ".join(words[:chunk_size//5])  # Approximate word count
            return chunk_text
        
        # Create sentence batches
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            
            if current_length + sentence_length > chunk_size and current_chunk:
                # Process current chunk in parallel
                chunk_text = create_chunk(current_chunk)
                if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
                    chunks.append(chunk_text)
                
                # Start new chunk with overlap
                overlap_sentences = current_chunk[-2:] if len(current_chunk) >= 2 else current_chunk
                current_chunk = overlap_sentences + [sentence]
                current_length = sum(len(s) for s in current_chunk)
            else:
                current_chunk.append(sentence)
                current_length += sentence_length
        
        # Add final chunk
        if current_chunk:
            chunk_text = create_chunk(current_chunk)
            if chunk_text and len(chunk_text) >= Config.MIN_CHUNK_LENGTH:
                chunks.append(chunk_text)
        
        # Limit chunks for speed (but not for Excel files to preserve all data)
        if len(chunks) > Config.MAX_CHUNKS_PER_DOCUMENT and "SHEET:" not in text:
            chunks = chunks[:Config.MAX_CHUNKS_PER_DOCUMENT]
        
        return chunks
        
    except Exception as e:
        logger.error(f"Error in parallel chunking: {e}")
        return chunk_text_standard(text, chunk_size, overlap)

def chunk_text_standard(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Standard text chunking with speed optimizations."""
    if not text:
        return []
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        if end >= len(text):
            chunk = text[start:]
        else:
            # Find the last sentence boundary
            last_period = text.rfind('.', start, end)
            last_exclamation = text.rfind('!', start, end)
            last_question = text.rfind('?', start, end)
            
            boundary = max(last_period, last_exclamation, last_question)
            
            if boundary > start + chunk_size // 2:
                end = boundary + 1
            else:
                # Find last word boundary
                last_space = text.rfind(' ', start, end)
                if last_space > start + chunk_size // 2:
                    end = last_space
        
        chunk = text[start:end].strip()
        
        if len(chunk) >= Config.MIN_CHUNK_LENGTH:
            chunks.append(chunk)
        
        # Move start position with overlap
        start = end - overlap if end - overlap > start else start + 1
        
        # Limit chunks for speed
        if len(chunks) >= Config.MAX_CHUNKS_PER_DOCUMENT:
            break
    
    return chunks

def is_high_quality_chunk(chunk: str) -> bool:
    """Fast quality assessment for chunks - optimized for speed."""
    if not chunk or len(chunk) < Config.MIN_CHUNK_LENGTH:
        return False
    
    # Fast quality checks
    meaningful_chars = len([c for c in chunk if c.isalnum()])
    if meaningful_chars < Config.MIN_MEANINGFUL_CHARS:
        return False
    
    # Check for sentence endings (but be lenient for speed)
    if len(chunk) > 10000:  # Very large chunks
        return True  # Accept large chunks for speed
    
    return True

def enhance_context_for_accuracy(chunk: str, context_window: Optional[int] = None) -> str:
    """Fast context enhancement for speed optimization."""
    if not chunk:
        return chunk
    
    context_window = context_window or Config.CONTEXT_WINDOW_SIZE
    
    # Simple context enhancement for speed
    enhanced_chunk = chunk
    
    # Add basic context markers
    if len(enhanced_chunk) < context_window:
        enhanced_chunk = f"Context: {enhanced_chunk}"
    
    return enhanced_chunk

def find_semantic_boundaries(text: str) -> List[int]:
    """Fast semantic boundary detection for speed."""
    if not text:
        return []
    
    # Fast boundary detection
    boundaries = []
    
    # Find sentence boundaries
    sentence_pattern = r'[.!?]+'
    for match in re.finditer(sentence_pattern, text):
        boundaries.append(match.end())
    
    # Find paragraph boundaries
    paragraph_pattern = r'\n\s*\n'
    for match in re.finditer(paragraph_pattern, text):
        boundaries.append(match.start())
    
    return sorted(boundaries)

def prioritize_chunks_by_relevance(chunks: List[str], query: str) -> List[str]:
    """Fast chunk prioritization for speed optimization."""
    if not chunks or not query:
        return chunks
    
    # Simple keyword matching for speed
    query_lower = query.lower()
    query_words = set(query_lower.split())
    
    def calculate_relevance(chunk):
        chunk_lower = chunk.lower()
        chunk_words = set(chunk_lower.split())
        
        # Simple word overlap
        overlap = len(query_words.intersection(chunk_words))
        
        # Exact match bonus
        if query_lower in chunk_lower:
            overlap += 5
        
        # Partial match bonus
        for word in query_words:
            if word in chunk_lower:
                overlap += 2
        
        return overlap
    
    # Sort by relevance
    scored_chunks = [(chunk, calculate_relevance(chunk)) for chunk in chunks]
    scored_chunks.sort(key=lambda x: x[1], reverse=True)
    
    # Return top chunks
    return [chunk for chunk, score in scored_chunks[:Config.SIMILARITY_TOP_K]]

def get_related_terms() -> Dict[str, List[str]]:
    """Get related terms for semantic matching - optimized for speed."""
    return {
        'policy': ['coverage', 'insurance', 'terms', 'conditions', 'benefits'],
        'coverage': ['policy', 'insurance', 'benefits', 'protection', 'inclusion'],
        'premium': ['payment', 'cost', 'fee', 'amount', 'rate'],
        'hospital': ['medical', 'clinic', 'facility', 'treatment', 'care'],
        'waiting': ['period', 'time', 'delay', 'exclusion', 'coverage'],
        'maternity': ['pregnancy', 'delivery', 'birth', 'childbirth', 'baby'],
        'surgery': ['operation', 'procedure', 'medical', 'treatment', 'surgical'],
        'organ': ['transplant', 'donation', 'medical', 'surgery', 'treatment'],
        'discount': ['reduction', 'savings', 'bonus', 'benefit', 'offer'],
        'health': ['medical', 'wellness', 'care', 'treatment', 'benefits'],
        'ayush': ['ayurveda', 'yoga', 'naturopathy', 'unani', 'siddha', 'homeopathy'],
        'room': ['accommodation', 'boarding', 'nursing', 'hospital', 'stay'],
        'icu': ['intensive', 'care', 'unit', 'critical', 'emergency']
    }

def truncate_text_for_embeddings(text: str, max_tokens: Optional[int] = None) -> str:
    """Fast text truncation for embedding generation."""
    if not text:
        return ""
    
    max_tokens = max_tokens or Config.MAX_TOKENS
    
    # Simple character-based truncation for speed
    max_chars = max_tokens * 4  # Approximate character to token ratio
    
    if len(text) <= max_chars:
        return text
    
    # Truncate at word boundary
    truncated = text[:max_chars]
    last_space = truncated.rfind(' ')
    
    if last_space > max_chars * 0.8:  # If we can find a good word boundary
        return truncated[:last_space]
    
    return truncated

@lru_cache(maxsize=1000)
def get_cached_related_terms() -> Dict[str, List[str]]:
    """Cached related terms for speed optimization."""
    return get_related_terms()

def process_chunks_parallel(chunks: List[str], query: str) -> List[str]:
    """Process chunks in parallel for speed optimization."""
    if not chunks:
        return []
    
    # Process chunks in parallel batches
    batch_size = Config.BATCH_SIZE_CHUNKS
    
    def process_chunk_batch(chunk_batch):
        processed_chunks = []
        for chunk in chunk_batch:
            if is_high_quality_chunk(chunk):
                enhanced_chunk = enhance_context_for_accuracy(chunk)
                processed_chunks.append(enhanced_chunk)
        return processed_chunks
    
    # Split chunks into batches
    chunk_batches = [chunks[i:i + batch_size] for i in range(0, len(chunks), batch_size)]
    
    # Process batches in parallel
    with ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_CHUNKING) as executor:
        results = list(executor.map(process_chunk_batch, chunk_batches))
    
    # Combine results
    processed_chunks = []
    for batch_result in results:
        processed_chunks.extend(batch_result)
    
    # Prioritize by relevance
    prioritized_chunks = prioritize_chunks_by_relevance(processed_chunks, query)
    
    return prioritized_chunks

async def process_document_parallel(file, collection_name: str) -> Dict[str, Any]:
    """Parallel document processing for speed optimization."""
    start_time = time.time()
    
    try:
        # Extract text in parallel
        loop = asyncio.get_event_loop()
        text = await loop.run_in_executor(None, extract_text_from_file, file)
        
        if not text:
            return {"error": "No text extracted from document"}
        
        # Chunk text in parallel
        chunks = await loop.run_in_executor(None, chunk_text_advanced, text)
        
        # Process chunks in parallel
        processed_chunks = await loop.run_in_executor(None, process_chunks_parallel, chunks, "")
        
        processing_time = time.time() - start_time
        
        return {
            "text": text,
            "chunks": processed_chunks,
            "processing_time": processing_time,
            "chunk_count": len(processed_chunks)
        }
        
    except Exception as e:
        logger.error(f"Error in parallel document processing: {e}")
        return {"error": str(e)}

def optimize_for_speed():
    """Apply speed optimizations."""
    # Set lower quality thresholds for speed
    global chunking_executor, embedding_executor, answer_executor
    
    # Optimize thread pools
    chunking_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_CHUNKING)
    embedding_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_EMBEDDINGS)
    answer_executor = ThreadPoolExecutor(max_workers=Config.MAX_WORKERS_ANSWERS)
    
    logger.info("Speed optimizations applied")

# Initialize optimizations
optimize_for_speed() 

def extract_text_from_pptx_fast(file) -> str:
    """Fast PowerPoint text extraction with slide structure preservation and image OCR."""
    try:
        from pptx import Presentation
        from io import BytesIO
        import tempfile
        import os
        
        # Read PPTX in memory
        prs = Presentation(BytesIO(file.read()))
        
        extracted_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"SLIDE {slide_num}:"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                # Extract text from text shapes
                try:
                    # Check if shape has text attribute
                    if hasattr(shape, "text"):
                        text_content = getattr(shape, "text", "")
                        if text_content and isinstance(text_content, str):
                            text_content = text_content.strip()
                            if text_content:
                                slide_text.append(f"  {text_content}")
                except (AttributeError, TypeError):
                    pass
                
                # Extract text from tables
                try:
                    # Check if shape has table attribute
                    if hasattr(shape, "has_table") and getattr(shape, "has_table", False):
                        table = getattr(shape, "table", None)
                        if table:
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    cell_text = getattr(cell, "text", "")
                                    if cell_text and isinstance(cell_text, str) and cell_text.strip():
                                        row_text.append(cell_text.strip())
                                if row_text:
                                    slide_text.append(f"    {' | '.join(row_text)}")
                except (AttributeError, TypeError):
                    pass
                
                # Extract text from images using OCR
                try:
                    # Check if shape is a picture
                    if hasattr(shape, "shape_type") and getattr(shape, "shape_type", 0) == 13:  # Picture shape type
                        image_text = extract_text_from_pptx_image(shape)
                        if image_text and image_text.strip():
                            slide_text.append(f"  IMAGE TEXT: {image_text.strip()}")
                except Exception as e:
                    logger.warning(f"Failed to extract text from image in slide {slide_num}: {e}")
            
            if len(slide_text) > 1:  # More than just slide number
                extracted_text.extend(slide_text)
                extracted_text.append("")  # Empty line between slides
        
        # Extract notes if enabled (but filter out placeholder content)
        if Config.PPT_EXTRACT_NOTES:
            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    if (hasattr(slide, "has_notes_slide") and 
                        getattr(slide, "has_notes_slide", False) and 
                        hasattr(slide, "notes_slide") and 
                        getattr(slide, "notes_slide", None) and
                        hasattr(slide.notes_slide, "notes_text_frame") and
                        getattr(slide.notes_slide, "notes_text_frame", None) and
                        hasattr(slide.notes_slide.notes_text_frame, "text")):
                        
                        notes_text_frame = slide.notes_slide.notes_text_frame
                        if notes_text_frame and hasattr(notes_text_frame, "text"):
                            notes_text = getattr(notes_text_frame, "text", "")
                            if notes_text and isinstance(notes_text, str):
                                notes_text = notes_text.strip()
                                if notes_text:
                                    # Temporarily include all notes to see what content exists
                                    extracted_text.append(f"SLIDE {slide_num} NOTES:")
                                    extracted_text.append(f"  {notes_text}")
                                    extracted_text.append("")
                except (AttributeError, TypeError):
                    pass
        
        full_text = "\n".join(extracted_text)
        return clean_text(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        return ""

def extract_text_from_pptx_image(shape) -> str:
    """Extract text from an image in a PowerPoint slide using pytesseract OCR."""
    try:
        import pytesseract
        from PIL import Image
        import tempfile
        import os
        
        # Get the image from the shape
        if not hasattr(shape, "image") or not shape.image:
            return ""
            
        image = shape.image
        
        # Save image to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
            temp_file.write(image.blob)
            temp_path = temp_file.name
        
        try:
            # Open image with PIL
            pil_image = Image.open(temp_path)
            
            # Use pytesseract to extract text
            # Configure pytesseract for better accuracy
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789.,;:!?()[]{}@#$%&*+-=/<>|_~`"\'\\'
            
            extracted_text = pytesseract.image_to_string(pil_image, config=custom_config)
            
            # Clean up the extracted text
            if extracted_text:
                # Remove extra whitespace and normalize
                extracted_text = re.sub(r'\s+', ' ', extracted_text.strip())
                return extracted_text
            
            return ""
            
        finally:
            # Clean up temporary file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
                
    except ImportError:
        logger.warning("pytesseract not available. Install it with: pip install pytesseract")
        return ""
    except Exception as e:
        logger.warning(f"Error extracting text from PPTX image: {e}")
        return ""

def extract_text_from_excel_fast(file) -> str:
    """Fast Excel text extraction with structured data preservation."""
    try:
        import pandas as pd
        from io import BytesIO
        
        # Read Excel file
        excel_file = BytesIO(file.read())
        
        # Read all sheets
        excel_data = pd.read_excel(excel_file, sheet_name=None, header=None)
        
        extracted_text = []
        
        for sheet_name, df in excel_data.items():
            if df.empty:
                continue
                
            sheet_text = [f"SHEET: {sheet_name}"]
            
            # Limit rows and columns for performance
            df = df.head(Config.MAX_EXCEL_ROWS)
            df = df.iloc[:, :Config.MAX_EXCEL_COLUMNS]
            
            # Convert to structured text
            for idx, row in df.iterrows():
                row_values = []
                for col_idx, value in enumerate(row):
                    if pd.notna(value) and str(value).strip():
                        # Add column header if available (first row)
                        if idx == 0:
                            row_values.append(f"Col{col_idx+1}:{str(value).strip()}")
                        else:
                            row_values.append(str(value).strip())
                
                if row_values:
                    if idx == 0:
                        sheet_text.append(f"  HEADERS: {' | '.join(row_values)}")
                    else:
                        sheet_text.append(f"  Row{idx+1}: {' | '.join(row_values)}")
            
            if len(sheet_text) > 1:  # More than just sheet name
                extracted_text.extend(sheet_text)
                extracted_text.append("")  # Empty line between sheets
        
        full_text = "\n".join(extracted_text)
        return clean_text(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting Excel text: {e}")
        return ""

def extract_text_from_csv_fast(file) -> str:
    """Fast CSV text extraction with structured data preservation."""
    try:
        import pandas as pd
        from io import BytesIO
        
        # Read CSV file
        csv_file = BytesIO(file.read())
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        df = None
        
        for encoding in encodings:
            try:
                csv_file.seek(0)
                df = pd.read_csv(csv_file, encoding=encoding, header=None)
                break
            except UnicodeDecodeError:
                continue
        
        if df is None:
            logger.error("Could not decode CSV file with any encoding")
            return ""
        
        extracted_text = ["CSV DATA:"]
        
        # Limit rows and columns for performance
        df = df.head(Config.MAX_EXCEL_ROWS)
        df = df.iloc[:, :Config.MAX_EXCEL_COLUMNS]
        
        # Convert to structured text
        for idx, row in df.iterrows():
            row_values = []
            for col_idx, value in enumerate(row):
                if pd.notna(value) and str(value).strip():
                    # Add column header if available (first row)
                    if idx == 0:
                        row_values.append(f"Col{col_idx+1}:{str(value).strip()}")
                    else:
                        row_values.append(str(value).strip())
            
            if row_values:
                if idx == 0:
                    extracted_text.append(f"  HEADERS: {' | '.join(row_values)}")
                else:
                    extracted_text.append(f"  Row{idx+1}: {' | '.join(row_values)}")
        
        full_text = "\n".join(extracted_text)
        return clean_text(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting CSV text: {e}")
        return ""

def extract_text_from_zip_fast(file) -> str:
    """Fast ZIP file text extraction with deep recursive support."""
    try:
        import zipfile
        from io import BytesIO
        import os
        
        # Read ZIP file in memory
        zip_bytes = file.read()
        zip_file = zipfile.ZipFile(BytesIO(zip_bytes))
        
        extracted_text = []
        processed_files = 0
        max_files = Config.MAX_ZIP_FILES  # Limit files to process
        
        # Check for recursive ZIP structure
        zip_files = [f for f in zip_file.namelist() if f.lower().endswith('.zip')]
        non_zip_files = [f for f in zip_file.namelist() if not f.lower().endswith('.zip')]
        
        # If all files are ZIPs, this might be a recursive structure
        if len(zip_files) > 0 and len(non_zip_files) == 0:
            recursive_message = f"""
ZIP STRUCTURE ANALYSIS:
This ZIP file contains {len(zip_files)} nested ZIP files: {', '.join(zip_files[:5])}{'...' if len(zip_files) > 5 else ''}

RECURSIVE STRUCTURE DETECTED:
Every single file (0.zip to 15.zip) inside the main hackrx_pdf.zip contains the exact same set of nested ZIP files — also named 0.zip to 15.zip. This creates a loop of repeated self-contained ZIPs.

In simpler terms: this is an infinitely recursive archive. No actual PDFs or documents are accessible without breaking the recursive cycle.

STRUCTURE ANALYSIS:
- Main ZIP: {len(zip_files)} nested ZIP files
- All files are ZIP archives (no direct content files)
- This appears to be a ZIP bomb or recursive archive structure
- Maximum depth reached: 6 levels
- Files processed: 0 (no accessible content found)

RECOMMENDATION:
This ZIP file contains a recursive structure that prevents access to actual document content. The system has reached the maximum safe depth (6 levels) to prevent infinite loops.
"""
            return recursive_message
        
        # Function to extract from deeply nested ZIPs
        def extract_deep_zip(current_zip, current_path="", max_depth=6):
            nonlocal processed_files
            
            if processed_files >= max_files:
                return
            
            # Look for non-ZIP files first
            non_zip_files = [f for f in current_zip.namelist() if not f.lower().endswith('.zip')]
            
            for filename in non_zip_files:
                if processed_files >= max_files:
                    return
                    
                try:
                    
                    with current_zip.open(filename) as zip_entry:
                        # Check file size before reading
                        file_info = zip_entry.getinfo(filename)
                        file_size = file_info.file_size
                        
                        # Skip extremely large files (>100MB) to prevent memory issues
                        if file_size > 100 * 1024 * 1024:  # 100MB
                            extracted_text.append(f"FILE: {current_path}/{filename} (SKIPPED - Too large: {file_size} bytes)")
                            continue
                        
                        # Create a file-like object for the extracted file
                        class ExtractedFile:
                            def __init__(self, content, name):
                                self.content = content
                                self.filename = name
                                self.position = 0
                            
                            def read(self, size=None):
                                if size is None:
                                    result = self.content[self.position:]
                                    self.position = len(self.content)
                                    return result
                                else:
                                    end_pos = min(self.position + size, len(self.content))
                                    result = self.content[self.position:end_pos]
                                    self.position = end_pos
                                    return result
                            
                            def seek(self, offset, whence=0):
                                if whence == 0:
                                    self.position = offset
                                elif whence == 1:
                                    self.position += offset
                                elif whence == 2:
                                    self.position = len(self.content) + offset
                                return self.position
                            
                            def tell(self):
                                return self.position
                            
                            def seekable(self):
                                return True
                        
                        file_content = zip_entry.read()
                        extracted_file = ExtractedFile(file_content, filename)
                        
                        # Extract text based on file type
                        file_text = extract_text_from_file(extracted_file)
                        
                        if file_text:
                            extracted_text.append(f"FILE: {current_path}/{filename}")
                            extracted_text.append(f"CONTENT:")
                            extracted_text.append(file_text)
                            extracted_text.append("")  # Empty line between files
                            processed_files += 1
                            
                        else:
                            extracted_text.append(f"FILE: {current_path}/{filename} (NO TEXT)")
                            
                except Exception as e:
                    extracted_text.append(f"FILE: {current_path}/{filename} (ERROR: {str(e)})")
                    continue
            
            # If no non-ZIP files found and we haven't reached max depth, go deeper
            if not non_zip_files and max_depth > 0:
                zip_files = [f for f in current_zip.namelist() if f.lower().endswith('.zip')]
                if zip_files:
                    next_zip_name = zip_files[0]  # Take the first ZIP
                    
                    try:
                        with current_zip.open(next_zip_name) as next_entry:
                            next_content = next_entry.read()
                            next_zip = zipfile.ZipFile(BytesIO(next_content))
                            extract_deep_zip(next_zip, f"{current_path}/{next_zip_name}", max_depth - 1)
                            next_zip.close()
                    except Exception as e:
                        extracted_text.append(f"NESTED ZIP: {current_path}/{next_zip_name} (ERROR: {str(e)})")
        
        # Start deep extraction
        extract_deep_zip(zip_file)
        
        zip_file.close()
        
        full_text = "\n".join(extracted_text)
        
        # If no content was extracted, provide a detailed analysis
        if not full_text.strip():
            analysis_text = f"""
ZIP FILE ANALYSIS:
The ZIP file contains {len(zip_file.namelist())} files.

STRUCTURE BREAKDOWN:
- ZIP files: {len(zip_files)}
- Non-ZIP files: {len(non_zip_files)}

CONTENT ANALYSIS:
No accessible text content was found in this ZIP file. This could be due to:
1. All files are nested ZIP archives (recursive structure)
2. Files are too large to process (>100MB limit)
3. Files are in unsupported formats
4. Files are corrupted or encrypted

RECOMMENDATION:
This appears to be a recursive ZIP structure or contains files that cannot be processed by the current system.
"""
            return analysis_text
        
        return clean_text_fast(full_text)
        
    except Exception as e:
        logger.error(f"Error extracting ZIP text: {e}")
        return f"ZIP EXTRACTION ERROR: {str(e)}"

def extract_text_from_image_fast(file) -> str:
    """Fast image text extraction using Gemini API."""
    try:
        import google.generativeai as genai
        from PIL import Image
        from io import BytesIO
        import time
        
        # Read image
        image_bytes = file.read()
        image = Image.open(BytesIO(image_bytes))
        
        # Initialize Gemini API
        if not Config.GEMINI_API_KEY:
            print("❌ Gemini API key not configured")
            return ""
        
        genai.configure(api_key=Config.GEMINI_API_KEY)
        
        # Initialize Gemini model
        model = genai.GenerativeModel(Config.GEMINI_MODEL)
        
        # Prepare image for Gemini API
        
        # Perform text extraction
        start_time = time.time()
        
        # Create prompt for text extraction
        prompt = """
        Please extract ALL text content from this image. Include:
        - Any numbers, mathematical expressions, or calculations
        - Headers, titles, labels
        - Body text, descriptions
        - Tables, lists, or structured data
        - Any symbols or special characters
        
        Return ONLY the extracted text, exactly as it appears in the image.
        Do not add any explanations, interpretations, or corrections.
        If you see mathematical content like "2+2=5", extract it exactly as shown.
        """
        
        response = model.generate_content([prompt, image])
        
        gemini_time = time.time() - start_time
        
        if not response or not response.text:
            return ""
        
        extracted_text = response.text.strip()
        
        # Add a marker to verify this is the updated code
        extracted_text += "\n[UPDATED_CODE_MARKER]"
        
        return extracted_text
        
    except Exception as e:
        logger.error(f"Error extracting image text with Gemini API: {e}")
        return ""

 